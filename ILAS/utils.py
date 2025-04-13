from django.conf import settings
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import subprocess
import os
import uuid
import glob

from .models import (
    InspectionAssignment,
    EstablishmentRegister,
    EstablishmentLicence,
    Establishment,
)


def inspector_assignments(user):
    """
    Retrieves the count of pending inspection assignments for a given inspector.

    Args:
        user (User): The inspector whose assignments are to be retrieved.

    Returns:
        int: Number of pending assignments (0 if none found)
        None: If an error occurs
    """
    try:
        assignments = InspectionAssignment.objects.filter(
            inspector=user, status="pending"
        ).count()
        return assignments if assignments else 0
    except Exception as e:
        return None


def mark_inspection_as_done(establishment):
    """
    Marks an inspection as completed in the database.

    Args:
        establishment: The establishment whose inspection is being marked complete

    Returns:
        None
    """
    try:
        inspection = InspectionAssignment.objects.get(establishment_id=establishment.id)
        inspection.status = "completed"
        inspection.save()
    except Exception as e:
        print(f"Error marking inspection as done: {e}")


def get_establishment_obj_by_register(register_id):
    """
    Retrieves an establishment object using its register ID.

    Args:
        register_id (int): The register ID associated with the establishment.

    Returns:
        Establishment: The establishment object if found
        None: If not found or error occurs
    """
    try:
        establishment = EstablishmentRegister.objects.get(id=register_id).establishment
        return establishment
    except Exception as e:
        print(f"Error getting establishment by register id: {e}")
        return None


def delete_files_except(directory_path, exception_file):
    """
    Delete all files in the given directory except for the specified exception file.

    Args:
        directory_path (str): Path to the directory where files should be deleted
        exception_file (str): Name of the file to be excluded from deletion
    """
    files_to_delete = glob.glob(os.path.join(directory_path, "*"))
    files_to_delete = [
        f for f in files_to_delete if os.path.basename(f) != exception_file
    ]

    for file in files_to_delete:
        try:
            os.remove(file)
            print(f"Deleted: {file}")
        except Exception as e:
            print(f"Error deleting {file}: {e}")


class PowerPointUpdater:
    """Handles updating and formatting PowerPoint presentations"""

    def __init__(self, presentation_path, cell_values):
        """
        Initialize the PowerPointUpdater.

        Args:
            presentation_path (str): Path to the PowerPoint template
            cell_values (dict): Dictionary mapping placeholder text to actual values
        """
        self.prs = Presentation(presentation_path)
        self.cell_values = cell_values

    def update_table_cells(self):
        """Update the text in table cells and shapes based on the cell_values dictionary"""
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    self._update_table(shape.table)
                elif hasattr(shape, "text"):
                    self._update_shape_text(shape)

    def _update_table(self, table):
        """Update text in table cells and apply formatting"""
        for row in table.rows:
            for cell in row.cells:
                if cell.text in self.cell_values:
                    cell.text = self.cell_values[cell.text]
                    self._center_align_text(cell)

    def _update_shape_text(self, shape):
        """Update text in non-table shapes"""
        if shape.text in self.cell_values:
            shape.text = self.cell_values[shape.text]

    def _center_align_text(self, cell):
        """Center-align text both vertically and horizontally in a table cell"""
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

    def save_presentation(self, output_path):
        """Save the updated presentation"""
        self.prs.save(output_path)


class PPTtoPDFConverter:
    """Handles conversion of PowerPoint files to PDF format"""

    LIBREOFFICE_PATH = "/usr/bin/libreoffice"

    @staticmethod
    def convert(input_file):
        """
        Convert a PowerPoint presentation to PDF using LibreOffice.

        Args:
            input_file (str): Path to the PowerPoint file

        Returns:
            str: Absolute path to the generated PDF file
        """
        output_dir = os.path.dirname(input_file)
        command = [
            PPTtoPDFConverter.LIBREOFFICE_PATH,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            input_file,
        ]
        subprocess.run(command, check=True)
        return os.path.abspath(input_file)


def create_license_report(licence_: EstablishmentLicence, establishment: Establishment, register):
    """
    Create a license report by generating a PDF from a PowerPoint template.

    Args:
        licence_: License details
        establishment: Establishment details
        register: Registration details

    Returns:
        str: Path to the generated PDF file
    """
    reports_temp_dir = os.path.join(settings.BASE_DIR, "static/files/reports_temp")
    template_path = os.path.join(reports_temp_dir, "license_report.pptx")
    output_pptx = os.path.join(reports_temp_dir, f"{uuid.uuid4()}.pptx")

    # Prepare data for the report
    cell_values = {
        "Owner_name": establishment.owner_name,
        "Register_number": str(register.id),
        "Establishment_name": str(establishment.establishment_name),
        "Id_number": str(establishment.owner_number),
        "License_category": str(licence_.main_category_id),
        "Issue_date": str(licence_.creation_date),
        "Expired_date": str(licence_.expiration_date),
        "Activity": str(licence_.activity.ar_name),
        "Address": str(establishment.get_address()),
        "License_number": str(licence_.number),
        "Phone_number": str(establishment.phone_number),
        "Email": str(establishment.email),
    }

    # Clean up old files
    delete_files_except(reports_temp_dir, "license_report.pptx")

    # Generate and convert the report
    updater = PowerPointUpdater(template_path, cell_values)
    updater.update_table_cells()
    updater.save_presentation(output_pptx)

    pdf_path = PPTtoPDFConverter.convert(os.path.abspath(output_pptx))
    os.remove(output_pptx)

    return pdf_path.split(".pptx")[0] + ".pdf"
